import argparse
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import logging

logging.basicConfig(level=logging.INFO)

def parse_args():
    parser = argparse.ArgumentParser(description="Generate Business Capability Map PPTX")
    parser.add_argument('--fontSizeLevel1', type=int, required=True)
    parser.add_argument('--fontSizeLevel2', type=int, required=True)
    parser.add_argument('--colorFillLevel1', type=str, required=True)
    parser.add_argument('--colorFillLevel2', type=str, required=True)
    parser.add_argument('--textColorLevel1', type=str, required=True)
    parser.add_argument('--textColorLevel2', type=str, required=True)
    parser.add_argument('--borderColor', type=str, required=True)
    parser.add_argument('--widthLevel2', type=float, required=True)
    parser.add_argument('--heightLevel2', type=float, required=True)
    parser.add_argument('--excelPath', type=str, required=False, default=None)
    parser.add_argument('--outputPath', type=str, required=False, default=None)
    return parser.parse_args()

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_colored_box(slide, left, top, width, height, text, fill_color, border_color, border_width, font_size, bold, text_color, align_left_top=False):
    try:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        rgb_fill = RGBColor(*hex_to_rgb(fill_color))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb_fill
        rgb_border = RGBColor(*hex_to_rgb(border_color))
        shape.line.color.rgb = rgb_border
        shape.line.width = Pt(border_width)
        # Make border invisible if border_color is None or border_width is 0
        if not border_color or border_width == 0:
            shape.line.fill.background()
        shape.text = text
        text_frame = shape.text_frame
        for p in text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = RGBColor(*hex_to_rgb(text_color))
        if align_left_top:
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.text import MSO_ANCHOR
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        else:
            text_frame.paragraphs[0].alignment = 1  # center
        # Make corners a little more rounded
        try:
            shape.adjustments[0] = 0.03
        except Exception:
            pass
        return shape
    except Exception as e:
        logging.error(f"Error in add_colored_box: {e}")
        return None

def add_business_capabilities(slide, df, args, prs_width, prs_height):
    level1_col = 'ID_1'
    level2_col = 'ID_2'
    level3_col = 'ID_3'
    capability_col = 'LEVEL_1_CAPABILITY'
    subcapability_col = 'LEVEL_2_CAPABILITY'
    subsubcapability_col = 'LEVEL_3_CAPABILITY'
    numbering_col = 'FULL_ID'

    # Layout constants
    margin = Inches(0.04)  # minimal gap between L1 columns (about 3pt)
    l1_text_margin = Inches(0.15)
    l2_box_margin = Inches(0.1)
    l2_box_side_margin = Inches(0.15)
    l3_box_margin = Inches(0.1)
    l1_box_min_width = Inches(2.0)
    l2_box_min_width = Inches(1.4)
    l3_box_min_width = Inches(1.0)
    l1_box_min_height = Inches(0.8)
    l1_box_min_height = float(l1_box_min_height) if l1_box_min_height is not None else 0.0
    l2_box_min_height = Inches(0.5)
    l2_box_min_height = float(l2_box_min_height) if l2_box_min_height is not None else 0.0
    l3_box_min_height = Inches(0.35)
    l3_box_min_height = float(l3_box_min_height) if l3_box_min_height is not None else 0.0
    l1_box_text_height = Inches(0.7)
    l2_box_text_height = Inches(0.5)
    l3_box_text_height = Inches(0.4)
    box_padding = Inches(0.15)
    child_left_pad = Inches(0.15)
    child_top_pad = Inches(0.02)
    vertical_spacing = Inches(0.1)
    min_box_height = Inches(0.2)
    min_spacing = Inches(0.05)
    min_padding = Inches(0.07)

    # Ensure all margin and padding values are defined as numbers (never None)
    l1_box_text_height = float(l1_box_text_height) if l1_box_text_height is not None else 0.0
    l2_box_text_height = float(l2_box_text_height) if l2_box_text_height is not None else 0.0
    l3_box_text_height = float(l3_box_text_height) if l3_box_text_height is not None else 0.0
    l2_box_side_margin = float(l2_box_side_margin) if l2_box_side_margin is not None else 0.0
    l2_box_margin = float(l2_box_margin) if l2_box_margin is not None else 0.0
    l3_box_margin = float(l3_box_margin) if l3_box_margin is not None else 0.0
    box_padding = float(box_padding) if box_padding is not None else 0.0
    child_left_pad = float(child_left_pad) if child_left_pad is not None else 0.0
    child_top_pad = float(child_top_pad) if child_top_pad is not None else 0.0
    vertical_spacing = float(vertical_spacing) if vertical_spacing is not None else 0.0
    min_box_height = float(min_box_height)
    min_spacing = float(min_spacing)
    min_padding = float(min_padding)

    # Normalize DataFrame column names for robustness
    df.columns = [c.strip().upper().replace(' ', '_') for c in df.columns]

    # 1. Build tree: {L1_ID: {row, children: {L2_ID: {row, children: {L3_ID: {row}}}}}}
    tree = {}
    for _, row in df.iterrows():
        l1 = row[level1_col]
        l2 = row[level2_col] if level2_col in row and pd.notnull(row[level2_col]) else None
        l3 = row[level3_col] if level3_col in row and pd.notnull(row[level3_col]) else None
        l1_key = str(l1)
        l2_key = str(l2) if l2 is not None else None
        l3_key = str(l3) if l3 is not None else None
        if l1_key not in tree:
            tree[l1_key] = {'row': row, 'children': {}}
        if l2_key:
            if l2_key not in tree[l1_key]['children']:
                tree[l1_key]['children'][l2_key] = {'row': row, 'children': {}}
            # Allow multiple L3s with the same ID by storing a list
            if l3_key:
                l3_children = tree[l1_key]['children'][l2_key]['children']
                if l3_key not in l3_children:
                    l3_children[l3_key] = []
                l3_children[l3_key].append({'row': row})

    # 2. Recursively compute heights for each node
    def compute_natural_height(node, level):
        if level == 3:
            return (l3_box_text_height or 0) + 2*(box_padding or 0)
        elif level == 2:
            children = node.get('children', {}) or {}
            if not children:
                return (l2_box_text_height or 0) + 2*(box_padding or 0)
            h = 0
            for c in children.values():
                h += (compute_natural_height(c, 3) or 0) + (vertical_spacing or 0)
            h -= (vertical_spacing or 0) if children else 0
            return h + 2*(box_padding or 0) + (l2_box_text_height or 0)
        elif level == 1:
            children = node.get('children', {}) or {}
            if not children:
                return (l1_box_text_height or 0) + 2*(box_padding or 0)
            h = 0
            for c in children.values():
                h += (compute_natural_height(c, 2) or 0) + (vertical_spacing or 0)
            h -= (vertical_spacing or 0) if children else 0
            return h + 2*(box_padding or 0) + (l1_box_text_height or 0)

    # 3. Recursively draw boxes
    def draw_node_scaled(node, level, left, top, width, scaling):
        scaling_f = float(scaling) if scaling is not None else 1.0
        if level == 1:
            l1_row = node['row']
            l1_num = str(l1_row[level1_col])
            l1_text = f"{l1_num}. {l1_row[capability_col]}"
            l1_min_height = float(l1_box_min_height) if l1_box_min_height is not None else 0.0
            l1_natural = float(compute_natural_height(node, 1) or 0.0)
            height = max(l1_min_height, scaling_f * l1_natural)
            pad = max(min_padding or 0, scaling_f * (box_padding or 0))
            header_h = l1_box_text_height * 0.7
            add_colored_box(
                slide, left, top, width, height,
                l1_text, args.colorFillLevel1, args.borderColor, 2, args.fontSizeLevel1, True, args.textColorLevel1, align_left_top=True
            )
            y = top + header_h + max(min_padding or 0, scaling_f * (child_top_pad or 0))
            for l2_id, l2_node in node.get('children', {}).items():
                l2_height = draw_node_scaled(l2_node, 2, left + max(min_padding or 0, scaling_f * (child_left_pad or 0)), y, width - 2*max(min_padding or 0, scaling_f * (child_left_pad or 0)), scaling_f)
                y += l2_height + max(min_spacing or 0, scaling_f * (vertical_spacing or 0))
            return height
        elif level == 2:
            l2_row = node['row']
            l1_num = str(l2_row[level1_col])
            l2_num = str(l2_row[level2_col])
            l2_text = f"{l1_num}.{l2_num} {l2_row[subcapability_col]}"
            pad = max(min_padding or 0, scaling_f * (box_padding or 0))
            children = node.get('children', {}) or {}
            # --- TIGHT L2 HEIGHT LOGIC ---
            HEIGHT_L3_CM = 1.0
            PADDING_CM = 0.2
            TOP_PADDING_CM = 0.2
            BOTTOM_MARGIN_CM = 0.2
            num_l3 = 0
            if not isinstance(children, dict) or not children:
                num_l3 = 0
            else:
                num_l3 = sum(len(l) for l in children.values())
                if not isinstance(num_l3, int):
                    num_l3 = 0
            height_l3 = Cm(HEIGHT_L3_CM)
            padding = Cm(PADDING_CM)
            top_padding = Cm(TOP_PADDING_CM)
            bottom_margin = Cm(BOTTOM_MARGIN_CM)
            if num_l3 > 0:
                height = top_padding + num_l3 * (height_l3 + padding) - padding + bottom_margin
            else:
                height = l2_box_min_height
            height = float(height)
            header_h = l2_box_text_height * 0.7
            l2_shape = add_colored_box(
                slide, left, top, width, height,
                l2_text, args.colorFillLevel2, args.borderColor, 1.5, args.fontSizeLevel2 + 1, True, args.textColorLevel2, align_left_top=True
            )
            # Place L3s strictly inside L2
            l3_left = left + max(min_padding or 0, scaling_f * (child_left_pad or 0))
            l3_width = width - 2 * max(min_padding or 0, scaling_f * (child_left_pad or 0))
            y = top + top_padding
            for l3_id, l3_nodes in children.items():
                for l3_node in l3_nodes:
                    l3_height = float(height_l3)
                    # Place L3 relative to L2's top/left
                    draw_node_scaled(l3_node, 3, l3_left, y, l3_width, scaling_f)
                    y += l3_height + float(padding)
            return height
        elif level == 3:
            l3_row = node['row']
            l1_num = str(l3_row[level1_col])
            l2_num = str(l3_row[level2_col])
            l3_num = str(l3_row[level3_col])
            l3_text = f"{l1_num}.{l2_num}.{l3_num} {l3_row[subsubcapability_col]}"
            l3_min_height = float(l3_box_min_height) if l3_box_min_height is not None else 0.0
            l3_natural = float((l3_box_text_height if l3_box_text_height is not None else 0.0) + 2*(box_padding if box_padding is not None else 0.0))
            scaling_f = float(scaling_f) if scaling_f is not None else 1.0
            height = max(float(l3_min_height), float(scaling_f * l3_natural))
            add_colored_box(
                slide, left, top, width, height,
                l3_text, args.colorFillLevel2, args.borderColor, 1, args.fontSizeLevel2 - 2, False, args.textColorLevel2, align_left_top=True
            )
            return height

    # 4. Layout Level 1 boxes horizontally, fit all on slide
    unique_level1s = list(tree.keys())
    n_level1 = len(unique_level1s)
    total_width = prs_width - 2*margin - (n_level1-1)*margin
    l1_width = max(l1_box_min_width, total_width / n_level1)
    l1_width = l1_width if l1_width is not None else 0.0
    margin_f = float(margin) if margin is not None else 0.0
    prs_width_f = float(prs_width) if prs_width is not None else 0.0
    n_level1 = int(n_level1) if n_level1 is not None else 0
    if l1_width * n_level1 + margin_f * (n_level1 - 1) > prs_width_f - 2*margin_f:
        l1_width = (prs_width_f - 2*margin_f - margin_f*(n_level1-1)) / n_level1 if n_level1 > 0 else l1_width
    available_height = prs_height - 2*margin
    l1_lefts = [margin + i*(l1_width + margin) for i in range(n_level1)]
    # --- GLOBAL SCALING FIX ---
    # Compute max natural height across all L1s
    max_natural_h = 0.0
    for l1_id in unique_level1s:
        l1_node = tree[l1_id]
        natural_h = float(compute_natural_height(l1_node, 1) or 0.0)
        if natural_h > max_natural_h:
            max_natural_h = natural_h
    scaling = min(1.0, available_height / max_natural_h) if max_natural_h > 0 else 1.0
    for i, l1_id in enumerate(unique_level1s):
        l1_node = tree[l1_id]
        draw_node_scaled(l1_node, 1, l1_lefts[i], margin, l1_width, scaling)

def generate_from_dataframe(df, args, output_path=None):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_business_capabilities(slide, df, args, prs.slide_width, prs.slide_height)
    if output_path is None:
        output_dir = os.path.join(os.path.dirname(__file__), 'output')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, 'Business_Capability_Map.pptx')
    else:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    logging.info(f"Saved presentation to {output_path}")
    print(f"Saved presentation to {output_path}")
    return output_path

def main():
    args = parse_args()
    excel_path = args.excelPath or os.path.join(os.path.dirname(__file__), 'excel_data', 'bcm_test_source.xlsx')
    if not os.path.exists(excel_path):
        logging.error(f"Excel file not found: {excel_path}")
        return
    df = pd.read_excel(excel_path)
    output_path = args.outputPath
    if output_path is None:
        output_dir = os.path.join(os.path.dirname(__file__), 'output')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, 'Business_Capability_Map.pptx')
    generate_from_dataframe(df, args, output_path=output_path)

if __name__ == "__main__":
    main()
